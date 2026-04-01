from __future__ import annotations

from pathlib import Path

import docx
import pandas as pd
import pytesseract
from fastapi import UploadFile
from PIL import Image
from pypdf import PdfReader
from pptx import Presentation

UPLOAD_DIR = Path("storage/uploads")


def _score_human_readable_filename(value: str) -> int:
    score = 0
    score += sum(char.isalpha() for char in value) * 2
    score += sum("\u0400" <= char <= "\u04FF" for char in value) * 3
    score -= value.count("Ð") * 4
    score -= value.count("Ñ") * 4
    score -= value.count("\\x") * 6
    score -= sum(ord(char) < 32 for char in value) * 10
    return score


def normalize_filename(name: str | None) -> str:
    raw_name = Path(name or "uploaded_file").name.strip() or "uploaded_file"
    candidates = [raw_name]

    if "\\x" in raw_name:
        try:
            candidates.append(raw_name.encode("utf-8").decode("unicode_escape"))
        except UnicodeDecodeError:
            pass

    for candidate in list(candidates):
        try:
            candidates.append(candidate.encode("latin1").decode("utf-8"))
        except (UnicodeEncodeError, UnicodeDecodeError):
            pass

    best = max(dict.fromkeys(candidates), key=_score_human_readable_filename)
    return best.replace("\x00", "").strip() or "uploaded_file"


def save_file(file: UploadFile) -> str:
    UPLOAD_DIR.mkdir(parents=True, exist_ok=True)
    file_name = normalize_filename(file.filename)
    path = UPLOAD_DIR / file_name

    with path.open("wb") as destination:
        destination.write(file.file.read())

    return str(path)


def read_txt(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="ignore")


def read_pdf(path: Path) -> str:
    text: list[str] = []
    reader = PdfReader(path)

    for page in reader.pages:
        page_text = page.extract_text()
        if page_text:
            text.append(page_text)

    return "\n".join(text)


def read_docx(path: Path) -> str:
    document = docx.Document(path)
    text: list[str] = []

    for paragraph in document.paragraphs:
        text.append(paragraph.text)

    return "\n".join(text)


def read_pptx(path: Path) -> str:
    presentation = Presentation(path)
    text: list[str] = []

    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)

    return "\n".join(text)


def read_xlsx(path: Path) -> str:
    dataframe = pd.read_excel(path)
    return dataframe.to_string()


def read_image(path: Path) -> str:
    image = Image.open(path)
    return pytesseract.image_to_string(image)


def extract_text_from_file(path: Path) -> str:
    ext = path.suffix.lower().lstrip(".")

    try:
        if ext == "txt":
            return read_txt(path)
        if ext == "pdf":
            return read_pdf(path)
        if ext == "docx":
            return read_docx(path)
        if ext == "pptx":
            return read_pptx(path)
        if ext in {"xls", "xlsx"}:
            return read_xlsx(path)
        if ext in {"png", "jpg", "jpeg"}:
            return read_image(path)
        return ""
    except Exception:
        return ""


def load_knowledge_base_documents() -> list[tuple[str, str]]:
    documents: list[tuple[str, str]] = []

    if not UPLOAD_DIR.exists():
        return documents

    for path in sorted(UPLOAD_DIR.iterdir()):
        if not path.is_file():
            continue
        text = extract_text_from_file(path)
        if text:
            documents.append((normalize_filename(path.name), text))

    return documents


def remove_uploaded_file(path: Path) -> None:
    try:
        if path.exists() and path.is_file():
            path.unlink()
    except Exception:
        return
